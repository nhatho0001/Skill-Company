"""
Template script v2 để build báo cáo KPI tháng theo chuẩn Tôn Đông Á.

CÁCH DÙNG:
1. Copy script này
2. Điều chỉnh phần CONFIG + REPORT CONTENT
3. Run → file .pptx xuất ra /mnt/user-data/outputs/

QUY TRÌNH SKILL ĐÃ ÁP DỤNG:
- Bước 2a: Drop 16 cột nhiễu
- Bước 2b: KHÔNG đề cập tên người, KHÔNG ghi thời gian thực hiện trong slide
- Bước 2b-bis: Tô màu header (đỏ = pending, navy = done)
- Bước 2d: Priority score chọn CV quan trọng (3-5 mỗi section)
- Bước 3b: Số slide động theo content density, cap 15
- Cách A (edit template) — giữ nguyên brand identity TDA v2

Template v2 (May 2026): 15 slide, 20×11.25in, font Inter, palette đỏ-primary + navy.
"""
import shutil
import pandas as pd
from pathlib import Path
from copy import deepcopy
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData


# ============================================================
# CONFIG — chỉnh theo từng báo cáo
# ============================================================
SKILL_DIR = Path("/mnt/skills/user/tda-kpi-report-slides")
TEMPLATE  = SKILL_DIR / "assets/template/report-template.pptx"

INPUT_FILE = Path("/mnt/user-data/uploads/congviec_thongke.xlsx")
OUT_PATH   = Path("/mnt/user-data/outputs/BaoCao_CNTT_T10-2025.pptx")

# Kỳ báo cáo
PERIOD_CURRENT = "10/2025"
PERIOD_NEXT    = "11/2025"
DEPT_FULLNAME  = "PHÒNG CÔNG NGHỆ THÔNG TIN"
DEPT_CODE      = "CNTT"  # CNTT / KT / HCM / KD / QLK


# ============================================================
# DESIGN TOKENS v2
# ============================================================
RED_PRIMARY   = RGBColor(0xFF, 0x00, 0x00)
NAVY          = RGBColor(0x00, 0x00, 0x99)
BLUE_TBL_HDR  = RGBColor(0x44, 0x72, 0xC4)
GRAY_TBL_EVEN = RGBColor(0xD9, 0xDB, 0xE7)
GRAY_TBL_ODD  = RGBColor(0xEC, 0xED, 0xF4)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK         = RGBColor(0x21, 0x21, 0x21)
FONT_HEAD     = "Inter"
FONT_BODY     = "Inter"


# ============================================================
# BƯỚC 2a: ĐỌC + DROP CỘT NHIỄU
# ============================================================
COLUMNS_TO_DROP = [
    "Lý do hủy", "Ngày tạo", "Lý do bỏ hoàn tất", "Ngày hoàn tất",
    "Đơn vị", "Chủ trì", "Phòng ban", "Đơn vị.1", "Liên hệ",
    "Phối hợp", "Để biết", "Hạng mục", "Đvt",
    "Số Ticket liên quan", "Khối lượng", "Người tạo",
    "Người hủy", "Ngày hủy", "Người hoàn tất", "Ngày Ticket liên quan"
]

# LƯU Ý: Data export từ hệ thống nội bộ đã được pre-filter theo phòng ở phía nguồn
# → KHÔNG cần lọc lại theo prefix Loại công việc. Dùng toàn bộ rows.

df = pd.read_excel(INPUT_FILE)
existing_drops = [c for c in COLUMNS_TO_DROP if c in df.columns]
df = df.drop(columns=existing_drops)
print(f"✓ Drop {len(existing_drops)} cột nhiễu. Tổng: {len(df)} rows (data đã pre-filter theo P.{DEPT_CODE})")


# ============================================================
# BƯỚC 2d: PRIORITY SCORE
# ============================================================
df["Duration_hours"] = (df["Đến ngày"] - df["Từ ngày"]).dt.total_seconds() / 3600

def priority_score(row):
    score = 0
    if row.get("Quan trọng") == "Quan trọng": score += 100
    if pd.notna(row.get("Dự án")) and str(row.get("Dự án")).strip(): score += 50
    dur = row.get("Duration_hours", 0)
    if pd.notna(dur):
        if dur > 24:    score += 30
        elif dur > 8:   score += 15
    if row.get("Khẩn cấp") == "Khẩn cấp": score += 20
    return score

df["Score"] = df.apply(priority_score, axis=1)
done    = df[df["Đã hoàn tất"] == True]
pending = df[df["Đã hoàn tất"] == False]

stats = dict(
    total=len(df), done=len(done), pending_count=len(pending),
    done_pct=round(len(done)/max(len(df),1)*100, 1),
)


# ============================================================
# BƯỚC 3: BUILD REPORT CONTENT
# ⚠️ KHÔNG đề cập tên người. KHÔNG ghi thời gian thực hiện.
# Mỗi section chọn 3-5 CV có Score cao nhất.
# ============================================================

# Section A — Hạ tầng (5 item icon_rows)
A_ITEMS = [
    {"header": "Hệ thống Mạng & Server",
     "body": "Các site hoạt động ổn định. Đã cấu hình liên kết mail phê duyệt EBS và chuyển Backup app 8007 qua Sanbox APP để tối ưu tài nguyên.",
     "is_pending": False},
    {"header": "Backup & Lưu trữ",
     "body": "Hệ thống Backup được theo dõi hàng ngày, tình trạng hoạt động ổn định, đảm bảo khả năng phục hồi dữ liệu.",
     "is_pending": False},
    {"header": "Bảo trì & License",
     "body": "Đã hoàn tất trình Kế hoạch bảo trì 2025 (đã được duyệt). Đồng thời, hoàn tất chuẩn bị Tờ trình gia hạn License Microsoft sắp tới.",
     "is_pending": False},
    {"header": "Hệ thống Camera",
     "body": "Tiếp tục thống kê tổng thể hệ thống camera tại các site. Ghi nhận tại TDM đã kiểm tra 199/205 kênh hoạt động.",
     "is_pending": False},
    {"header": "Hỗ trợ Đặc biệt",
     "body": "Đã lên hợp đồng mua Chữ ký số VNPT SmartCA phục vụ cho các cơ quan chức năng (COCQ).",
     "is_pending": False},
]

# Section B — Timeline 4 bước (numbered_zigzag_4)
B_ITEMS = [
    {"header": "Vận hành & Bảo trì",
     "body": "Hoàn thành ký duyệt Hợp đồng bảo trì 06 tháng Database + DRSite, đảm bảo tính liên tục và an toàn dữ liệu cho hệ thống ERP cốt lõi.",
     "is_pending": False},
    {"header": "Hóa đơn Điện tử",
     "body": "Theo dõi ổn định sau Golive HĐĐT cải tiến ND70. Đã hoàn thành và ký duyệt các biên bản nghiệm thu liên quan đến dự án này.",
     "is_pending": False},
    {"header": "Nâng cấp ERP",
     "body": "Hợp đồng dự án Nâng cấp ERP đã được phê duyệt. Đang thực hiện theo kế hoạch và đã hoàn thành các Biên bản Khảo sát (BBKS) sơ lược ban đầu.",
     "is_pending": True},
    {"header": "QR Code Thép Hộp",
     "body": "Hợp đồng triển khai đã hoàn thành ký duyệt, chính thức bắt đầu dự án số hóa quản lý sản phẩm.",
     "is_pending": False},
]

# Section C — 3 cards (cards_3col)
C_ITEMS = [
    {"header": "EOFFICE & HCM",
     "body": "Đã thực hiện điều chỉnh phòng/ban và thay đổi chức danh trên hệ thống theo Cơ cấu tổ chức mới của Khối Sản Xuất và Phòng Quản lý kho.",
     "is_pending": False},
    {"header": "Ngừng Kho Trung Chuyển (ERP)",
     "body": "Thông báo ngưng sử dụng kho trung chuyển trên hệ thống ERP chính thức từ ngày 01/11/2025, tối ưu hóa quy trình quản lý kho hàng.",
     "is_pending": False},
    {"header": "Chiết Khấu Bán Hàng (CKBH)",
     "body": "Đã tổ chức họp trao đổi phạm vi phát sinh các Change Request (CR) sau Golive với đối tác iERP.",
     "is_pending": True},
]

# Pending Section — 4 ô đỏ
PENDING_ITEMS = [
    {"header": "Giám Sát Hạ Tầng Mạng",
     "body": "Tiếp tục giám sát tiến độ thi công mạng tại TDA.LA và ĐA.PM1, đảm bảo chất lượng và đúng tiến độ."},
    {"header": "Quản Lý Lưu Trữ Camera",
     "body": "Trình phương án quản lý và mở rộng thời gian lưu trữ camera theo quy định của Nghị định 167."},
    {"header": "Cập Nhật Bảo Mật (DLP)",
     "body": "Đánh giá và cập nhật giải pháp Data Loss Prevention (DLP) với các nền tảng như ManageEngine/MS 365/Trellix."},
    {"header": "Nâng Cấp ERP",
     "body": "Hoàn thành test các nghiệp vụ ERP/APP TDA/QRCODE/DMS"},
]

# Table data (cho slide 11 nếu giữ)
TABLE_HEADERS = ["Nhóm công việc", "Tổng", "Đã xong", "% Hoàn thành"]
TABLE_ROWS = [
    ["A. Hạ tầng CNTT",     45, 42, "93%"],
    ["B. Hệ thống ERP",     28, 25, "89%"],
    ["C. Eoffice / HCM",    18, 17, "94%"],
    ["D. CĐS & AI",         12,  9, "75%"],
    ["E. Tổng hợp",          6,  6, "100%"],
]

# Chart data (slide 13 = donut)
CHART13_CATS  = ["Hạ tầng", "ERP", "Eoffice", "CĐS/AI"]
CHART13_DATA  = {"% Hoàn thành": [42, 25, 17, 9]}


# ============================================================
# HELPERS (xem building-blocks.md để biết bản đầy đủ)
# ============================================================
def _iter_all_shapes(shapes_collection):
    """Iterate đệ quy qua tất cả shape, kể cả shape trong Group."""
    for shape in shapes_collection:
        yield shape
        if shape.shape_type == 6:  # GROUP
            yield from _iter_all_shapes(shape.shapes)


def set_text(tf, text, size=14, bold=False, color=BLACK,
             font=FONT_BODY, align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def replace_shape_by_name(slide, shape_name, new_text):
    for shape in _iter_all_shapes(slide.shapes):
        if shape.name == shape_name and shape.has_text_frame:
            tf = shape.text_frame
            p0 = tf.paragraphs[0]
            if p0.runs:
                p0.runs[0].text = new_text
                for r in p0.runs[1:]:
                    r.text = ""
            for p in tf.paragraphs[1:]:
                for r in p.runs:
                    r.text = ""
            return True
    return False


def replace_text_anywhere(slide, old, new):
    for shape in _iter_all_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            if old in full:
                new_full = full.replace(old, new)
                if para.runs:
                    para.runs[0].text = new_full
                    for r in para.runs[1:]:
                        r.text = ""


def replace_textframe_by_name(slide, shape_name, new_text):
    """Replace TOÀN BỘ text frame (gộp nhiều paragraph thành 1).
    Dùng khi text gốc trải qua nhiều paragraph (vd subtitle chart slide 13/14).
    """
    for shape in _iter_all_shapes(slide.shapes):
        if shape.name == shape_name and shape.has_text_frame:
            tf = shape.text_frame
            p0 = tf.paragraphs[0]
            if p0.runs:
                p0.runs[0].text = new_text
                for r in p0.runs[1:]:
                    r.text = ""
            else:
                run = p0.add_run()
                run.text = new_text
            for p in tf.paragraphs[1:]:
                for r in p.runs:
                    r.text = ""
            return True
    return False


def replace_cover_period(slide_cover, period_current, period_next, department=None):
    """Template v2: 6 paragraph riêng biệt — không cần <a:br/> hack như v1."""
    for shape in _iter_all_shapes(slide_cover.shapes):
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if "KẾT QUẢ" not in tf.text and "BÁO CÁO" not in tf.text:
            continue
        for para in tf.paragraphs:
            full = "".join(r.text for r in para.runs)
            if "KẾT QUẢ" in full and para.runs:
                para.runs[0].text = f"KẾT QUẢ THÁNG {period_current}"
                for r in para.runs[1:]:
                    r.text = ""
            elif ("VÀ KẾ HOẠCH" in full or "KẾ HOẠCH" in full) and para.runs:
                para.runs[0].text = f"VÀ KẾ HOẠCH THÁNG {period_next}"
                for r in para.runs[1:]:
                    r.text = ""
            elif department and "PHÒNG" in full and para.runs:
                para.runs[0].text = department
                for r in para.runs[1:]:
                    r.text = ""
        return True
    return False


def force_header_color(slide, shape_name, is_pending: bool):
    target = RED_PRIMARY if is_pending else NAVY
    for shape in _iter_all_shapes(slide.shapes):
        if shape.name != shape_name or not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    run.font.color.rgb = target
        return True
    return False


def shrink_title_if_long(slide, shape_name, max_chars=50, small_size=32):
    for shape in _iter_all_shapes(slide.shapes):
        if shape.name == shape_name and shape.has_text_frame:
            text = shape.text_frame.text
            if len(text) > max_chars:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(small_size)
                return True
    return False


def replace_table_in_slide(slide, headers, rows,
                            x=2.25, y=3.38, w=14.42, max_h=7.5):
    for shape in list(slide.shapes):
        if shape.has_table:
            sp = shape._element
            sp.getparent().remove(sp)
            break

    n_rows = len(rows) + 1
    n_cols = len(headers)
    h = min(0.7 + 0.5 * len(rows), max_h)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = tbl_shape.table

    for c, hd in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_TBL_HDR
        tf = cell.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = str(hd)
        run.font.bold = True; run.font.size = Pt(15)
        run.font.color.rgb = WHITE; run.font.name = FONT_HEAD
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for r, row_data in enumerate(rows, start=1):
        bg = GRAY_TBL_EVEN if r % 2 == 1 else GRAY_TBL_ODD
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            tf = cell.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = str(val)
            run.font.bold = True; run.font.size = Pt(14)
            run.font.color.rgb = BLACK; run.font.name = FONT_BODY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl_shape


def replace_chart_data(slide, categories, series_data):
    for shape in _iter_all_shapes(slide.shapes):
        if shape.has_chart:
            chart_data = CategoryChartData()
            chart_data.categories = categories
            for s_name, values in series_data.items():
                chart_data.add_series(s_name, values)
            shape.chart.replace_data(chart_data)
            return shape.chart
    return None


def delete_slide(prs, slide_idx):
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    rId = slides_list[slide_idx].rId
    prs.part.drop_rel(rId)
    xml_slides.remove(slides_list[slide_idx])


def debug_slide_shapes(prs, slide_idx=None):
    indices = [slide_idx] if slide_idx is not None else range(len(prs.slides))
    for i in indices:
        slide = prs.slides[i]
        print(f"\n=== Slide {i+1} ===")
        for shape in _iter_all_shapes(slide.shapes):
            preview = ""
            if shape.has_text_frame and shape.text_frame.text.strip():
                preview = " | " + shape.text_frame.text.replace("\n", " ⏎ ")[:80]
            elif shape.has_table:
                tbl = shape.table
                preview = f" | [TABLE {len(tbl.rows)}r x {len(tbl.columns)}c]"
            elif shape.has_chart:
                preview = f" | [CHART]"
            print(f"  [{shape.name}]{preview}")


# ============================================================
# BƯỚC 5: EDIT TEMPLATE
# ============================================================
shutil.copy(TEMPLATE, OUT_PATH)
prs = Presentation(str(OUT_PATH))

# Debug TRƯỚC khi sửa (uncomment khi test):
# debug_slide_shapes(prs)


# ---------- SLIDE 1: Cover ----------
replace_cover_period(prs.slides[0], PERIOD_CURRENT, PERIOD_NEXT, DEPT_FULLNAME)


# ---------- SLIDE 2: TOC ----------
s2 = prs.slides[1]
shrink_title_if_long(s2, "TextBox 3", max_chars=40, small_size=32)
replace_text_anywhere(s2,
    "Báo cáo này tập trung vào các kết quả đã đạt được và kế hoạch hành động chính của Phòng CNTT trong tháng 10 và định hướng cho tháng 11.",
    f"Báo cáo tổng hợp {stats['done']}/{stats['total']} công việc hoàn thành ({stats['done_pct']}%) và {stats['pending_count']} công việc tồn đọng của {DEPT_FULLNAME} trong tháng {PERIOD_CURRENT}, cùng định hướng tháng {PERIOD_NEXT}."
)
# 5 mục TOC: TextBox 11/12 = mục 1, TextBox 19/20 = mục 2, ... (chạy debug để confirm)


# ---------- SLIDE 3: Section A — icon_rows ----------
s3 = prs.slides[2]
# Mapping (verify bằng debug):
# Title = TextBox 3
# Item 1 header = TextBox 4, body = TextBox 5
# Item 2 header = TextBox 6, body = TextBox 7
# Item 3 header = TextBox 8, body = TextBox 9
# Item 4 header = TextBox 10, body = TextBox 11
# Item 5 header = TextBox 12, body = TextBox 13
HEADER_NAMES_S3 = ["TextBox 4", "TextBox 6", "TextBox 8", "TextBox 10", "TextBox 12"]
BODY_NAMES_S3   = ["TextBox 5", "TextBox 7", "TextBox 9", "TextBox 11", "TextBox 13"]

for i, item in enumerate(A_ITEMS[:5]):
    replace_shape_by_name(s3, HEADER_NAMES_S3[i], item["header"])
    replace_shape_by_name(s3, BODY_NAMES_S3[i], item["body"])
    force_header_color(s3, HEADER_NAMES_S3[i], item["is_pending"])


# ---------- SLIDE 4: Section B — numbered_zigzag_4 ----------
s4 = prs.slides[3]
# Item 1: TextBox 13 + 14, Item 2: TextBox 21 + 22
# Item 3: TextBox 29 + 30, Item 4: TextBox 37 + 38
S4_HEADERS = ["TextBox 13", "TextBox 21", "TextBox 29", "TextBox 37"]
S4_BODIES  = ["TextBox 14", "TextBox 22", "TextBox 30", "TextBox 38"]
for i, item in enumerate(B_ITEMS[:4]):
    replace_shape_by_name(s4, S4_HEADERS[i], item["header"])
    replace_shape_by_name(s4, S4_BODIES[i], item["body"])
    force_header_color(s4, S4_HEADERS[i], item["is_pending"])


# ---------- SLIDE 5: Section C — cards_3col ----------
s5 = prs.slides[4]
# 3 cards: TextBox 9+10, TextBox 15+16, TextBox 21+22
S5_HEADERS = ["TextBox 9", "TextBox 15", "TextBox 21"]
S5_BODIES  = ["TextBox 10", "TextBox 16", "TextBox 22"]
for i, item in enumerate(C_ITEMS[:3]):
    replace_shape_by_name(s5, S5_HEADERS[i], item["header"])
    replace_shape_by_name(s5, S5_BODIES[i], item["body"])
    force_header_color(s5, S5_HEADERS[i], item["is_pending"])


# ---------- SLIDE 6, 7: tuỳ data, có thể skip/replace ----------
# delete_slide(prs, 5)  # nếu không có nội dung CĐS/AI


# ---------- SLIDE 8: 3-project image card (XÓA nếu không cần) ----------
# Nếu không có dự án trọng điểm → xóa slide này
delete_slide_indices = []
# delete_slide_indices.append(7)  # slide 8 (0-indexed = 7)


# ---------- SLIDE 9: Pending ----------
s9 = prs.slides[8]
# Header thay theo PERIOD_NEXT
replace_text_anywhere(s9, "TỒN ĐỌNG & TRỌNG TÂM THÁNG 11/2025",
                      f"TỒN ĐỌNG & TRỌNG TÂM THÁNG {PERIOD_NEXT}")
# 4 items: TextBox 12+13, 21+22, 30+31, 39+40
S9_HEADERS = ["TextBox 12", "TextBox 21", "TextBox 30", "TextBox 39"]
S9_BODIES  = ["TextBox 13", "TextBox 22", "TextBox 31", "TextBox 40"]
for i, item in enumerate(PENDING_ITEMS[:4]):
    replace_shape_by_name(s9, S9_HEADERS[i], item["header"])
    replace_shape_by_name(s9, S9_BODIES[i], item["body"])
    # Pending → tô đỏ tất cả header
    force_header_color(s9, S9_HEADERS[i], is_pending=True)


# ---------- SLIDE 10: Timeline (placeholder generic) ----------
# Nếu không có lộ trình theo năm → xóa
delete_slide_indices.append(9)  # slide 10


# ---------- SLIDE 11: Table ----------
s11 = prs.slides[10]
replace_text_anywhere(s11, "Table page heading",
                      f"BẢNG TỔNG HỢP CÔNG VIỆC THÁNG {PERIOD_CURRENT}")
replace_text_anywhere(s11,
    "Subheading that introduces the table’s contents and highlights its key metrics.",
    f"Tỷ lệ hoàn thành theo nhóm công việc của {DEPT_FULLNAME}.")
replace_table_in_slide(s11, TABLE_HEADERS, TABLE_ROWS)


# ---------- SLIDE 12: 3-project image card (Hoạt động khác) ----------
# Có thể giữ nếu có 3 hoạt động trọng điểm khác, hoặc xóa
delete_slide_indices.append(11)  # slide 12


# ---------- SLIDE 13: Donut chart ----------
s13 = prs.slides[12]
# ⚠️ Title slide 13 panel hẹp ~4" → giữ ≤ 22-25 ký tự (Quirk 3)
replace_shape_by_name(s13, "TextBox 7", "PHÂN BỔ KẾT QUẢ T10")  # 19 ký tự, không wrap
replace_textframe_by_name(s13, "TextBox 8",
    "Tỷ trọng % hoàn thành theo nhóm CV chính.")
replace_textframe_by_name(s13, "TextBox 5",
    f"Dữ liệu tổng hợp tháng {PERIOD_CURRENT}, phản ánh tỷ trọng hoàn thành theo từng nhóm.")
replace_chart_data(s13, CHART13_CATS, CHART13_DATA)


# ---------- SLIDE 14: Column chart (XÓA nếu không cần chart 2) ----------
delete_slide_indices.append(13)  # slide 14


# ---------- XÓA TỪ CUỐI LÊN ----------
for idx in sorted(set(delete_slide_indices), reverse=True):
    delete_slide(prs, idx)


# ---------- Save ----------
prs.save(str(OUT_PATH))
print(f"\n✓ Đã lưu: {OUT_PATH}")
print(f"  Slides: {len(prs.slides)} | Done: {stats['done']}/{stats['total']} ({stats['done_pct']}%)")


# ============================================================
# BƯỚC 6: QA (chạy sau khi save)
# ============================================================
# 1. Text QA:
#    extract-text /mnt/user-data/outputs/<file>.pptx
#    extract-text /mnt/user-data/outputs/<file>.pptx | grep -iE "lorem|column 1|content|page heading|TODO|\bx{3,}\b"
# 2. Visual QA:
#    cd /home/claude
#    python /mnt/skills/public/pptx/scripts/office/soffice.py --headless --convert-to pdf <out_path>
#    rm -f slide-*.jpg
#    pdftoppm -jpeg -r 100 <pdf> slide
#    # Review từng slide-*.jpg
# 3. Fix lỗi overflow bằng cách rút gọn header (≤20 ký tự cho 3-col, ≤16 cho 4-col)
