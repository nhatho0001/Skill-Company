"""
Template script để build báo cáo KPI tháng theo chuẩn Tôn Đông Á.

CÁCH DÙNG:
1. Copy script này
2. Điều chỉnh phần DATA INPUT + REPORT CONTENT
3. Run → file .pptx xuất ra /mnt/user-data/outputs/

QUY TRÌNH SKILL ĐÃ ÁP DỤNG:
- Bước 2a: Drop 20 cột nhiễu
- Bước 2b: KHÔNG đề cập tên người trong slide
- Bước 2d: Priority score chọn CV quan trọng
- Cách A (edit template) — giữ nguyên brand identity
"""
import shutil
import pandas as pd
from pathlib import Path
from copy import deepcopy
from lxml import etree

from pptx import Presentation
from pptx.oxml.ns import qn

# ============================================================
# CONFIG — chỉnh theo từng báo cáo
# ============================================================
SKILL_DIR = Path("/home/claude/tda-kpi-report-slides")  # hoặc path skill thực
TEMPLATE  = SKILL_DIR / "assets/template/report-template.pptx"

# Thay path & tên file theo kỳ báo cáo
INPUT_FILE = Path("/mnt/user-data/uploads/congviec_thongke.xlsx")
OUT_PATH   = Path("/mnt/user-data/outputs/BaoCao_CNTT_T04-2026.pptx")

# Kỳ báo cáo
PERIOD_CURRENT = "04/2026"   # tháng hiện tại
PERIOD_NEXT    = "05/2026"   # tháng kế tiếp

# ============================================================
# BƯỚC 2a: ĐỌC + DROP CỘT NHIỄU (BẮT BUỘC)
# ============================================================
COLUMNS_TO_DROP = [
    "Lý do hủy", "Ngày tạo", "Lý do bỏ hoàn tất", "Ngày hoàn tất",
    "Đơn vị", "Chủ trì", "Phòng ban", "Đơn vị.1", "Liên hệ",
    "Phối hợp", "Để biết", "Hạng mục", "Đvt",
    "Số Ticket liên quan", "Khối lượng", "Người tạo",
    "Người hủy", "Ngày hủy", "Người hoàn tất", "Ngày Ticket liên quan"
]

df = pd.read_excel(INPUT_FILE)
existing_drops = [c for c in COLUMNS_TO_DROP if c in df.columns]
df = df.drop(columns=existing_drops)
print(f"✓ Drop {len(existing_drops)} cột nhiễu. Còn lại: {df.columns.tolist()}")

# ============================================================
# BƯỚC 2d: PRIORITY SCORE
# ============================================================
df["Duration_hours"] = (df["Đến ngày"] - df["Từ ngày"]).dt.total_seconds() / 3600

def priority_score(row):
    score = 0
    if row.get("Quan trọng") == "Quan trọng": score += 100
    if pd.notna(row.get("Dự án")) and str(row.get("Dự án")).strip(): score += 50
    dur = row.get("Duration_hours", 0)
    if pd.notna(dur):
        if dur > 24: score += 30
        elif dur > 8: score += 15
    if row.get("Khẩn cấp") == "Khẩn cấp": score += 20
    return score

df["Score"] = df.apply(priority_score, axis=1)

# ============================================================
# PHÂN LOẠI THEO NHÓM (tuỳ chỉnh theo dữ liệu của bạn)
# ============================================================
def classify(loai):
    if pd.isna(loai): return "E"
    s = str(loai)
    if "CNTT-HT" in s or "Firewall" in s or "Fileserver" in s: return "A"
    if "CNTT-ERP" in s or s.startswith("ERP"): return "B"
    if "EOffice" in s or "eOffice" in s or "HCM" in s or "Xử lý công việc quy trình" in s: return "C"
    return "E"

df["Nhom"] = df["Loại công việc"].apply(classify)
done = df[df["Đã hoàn tất"] == True]
pending = df[df["Đã hoàn tất"] == False]

stats = dict(
    total=len(df), done=len(done), pending_count=len(pending),
    done_pct=round(len(done)/len(df)*100, 1),
    A=len(done[done["Nhom"] == "A"]),
    B=len(done[done["Nhom"] == "B"]),
    C=len(done[done["Nhom"] == "C"]),
    E=len(done[done["Nhom"] == "E"]),
)

# ============================================================
# BƯỚC 3: BUILD REPORT CONTENT
# ============================================================
# ⚠️ QUAN TRỌNG: KHÔNG đề cập tên người chủ trì trong các đoạn text dưới đây!
# Báo cáo cấp phòng/ban chỉ nói về công việc, dự án, hệ thống, đơn vị.

# Phần này Claude cần suy luận từ df (top CV có Score cao nhất cho từng nhóm)
# Sau đây là ví dụ dựa trên data mẫu T04/2026

# Section A — các CV hạ tầng quan trọng nhất
A_ITEMS = [
    ("Header 1", "Body 1 — thời gian, trạng thái, nội dung chính..."),
    # ... 3-5 items
]

# Section B — 4 dự án ERP trọng điểm
B_ITEMS = [
    ("Header 1", "Body 1"),
    # ... 4 items (template có 4 bước)
]

# Section C — 3 quy trình trọng điểm
C_ITEMS = [
    ("Header 1", "Body 1"),
    # ... 3 items
]

# Section D (nếu giữ) — hoặc sẽ XÓA slide này
# KEEP_SLIDE_D = False

# Section E — báo cáo, hành chính
E_ITEMS = [
    ("Header 1", "Body 1"),
    # ... 4 items
]

# Pending — 4 CV tồn đọng ưu tiên
PENDING_ITEMS = [
    ("Header 1", "Body 1 — có icon ⭐ cho Quan trọng, 🔴 cho Khẩn cấp"),
    # ... 4 items
]

# Kiến nghị — 3 trọng tâm T kế
KIENNGHI_ITEMS = [
    ("Header 1", "Body 1"),
    # ... 3 items
]

# ============================================================
# BƯỚC 5: EDIT TEMPLATE
# ============================================================
shutil.copy(TEMPLATE, OUT_PATH)
prs = Presentation(str(OUT_PATH))


# ---------- Helper functions ----------
def replace_text_anywhere(slide, old, new):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            if old in full:
                new_full = full.replace(old, new)
                if para.runs:
                    para.runs[0].text = new_full
                    for r in para.runs[1:]:
                        r.text = ""

def replace_shape_by_name(slide, shape_name, new_text):
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            tf = shape.text_frame
            p0 = tf.paragraphs[0]
            if p0.runs:
                p0.runs[0].text = new_text
                for r in p0.runs[1:]:
                    r.text = ""
            for p in tf.paragraphs[1:]:
                for r in p.runs:
                    r.text = ""
            return True
    return False

def delete_slide(prs, slide_idx):
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    rId = slides_list[slide_idx].rId
    prs.part.drop_rel(rId)
    xml_slides.remove(slides_list[slide_idx])


# ---------- SLIDE 1: Cover ----------
# ⚠️ Quirk 2: multi-run paragraph — cần thêm <a:br/> để giữ line break
s1 = prs.slides[0]
for shape in s1.shapes:
    if shape.has_text_frame and "KẾT QUẢ" in shape.text_frame.text:
        for para in shape.text_frame.paragraphs:
            if "KẾT QUẢ" in "".join(r.text for r in para.runs):
                if para.runs:
                    para.runs[0].text = f"KẾT QUẢ THÁNG {PERIOD_CURRENT}"
                    for r in para.runs[1:]:
                        r.text = ""
                    p_xml = para._p
                    etree.SubElement(p_xml, qn("a:br"))
                    new_r = deepcopy(para.runs[0]._r)
                    new_r.find(qn("a:t")).text = f"VÀ KẾ HOẠCH THÁNG {PERIOD_NEXT}"
                    p_xml.append(new_r)
                break

# ---------- SLIDE 2: Mục lục ----------
s2 = prs.slides[1]
replace_text_anywhere(s2,
    "Báo cáo này tập trung vào các kết quả đã đạt được và kế hoạch hành động chính của Phòng CNTT trong tháng 10 và định hướng cho tháng 11.",
    f"Báo cáo tổng hợp {stats['done']} công việc đã hoàn tất và {stats['pending_count']} công việc tồn đọng của Phòng CNTT trong tháng {PERIOD_CURRENT}, cùng định hướng cho tháng {PERIOD_NEXT}."
)
# TODO: replace header và description của 5 mục trong TOC

# ---------- SLIDE 3: Section A ----------
s3 = prs.slides[2]
# Shape names: Text 1, Text 3, Text 5, Text 7, Text 9 = headers
#              Text 2, Text 4, Text 6, Text 8, Text 10 = bodies
for i, (header, body) in enumerate(A_ITEMS):
    replace_shape_by_name(s3, f"Text {1 + 2*i}", header)
    replace_shape_by_name(s3, f"Text {2 + 2*i}", body)

# ---------- SLIDE 4, 5, 7 (E), 8 (Pending), 9 (Kiến nghị) ----------
# Xem từng slide: chạy debug script trước để lấy mapping shape_name
# TODO

# ---------- XÓA SLIDE D (6) nếu không có nội dung CĐS/AI ----------
# delete_slide(prs, 5)

# ---------- Save ----------
prs.save(str(OUT_PATH))
print(f"✓ Đã lưu: {OUT_PATH}")
print(f"  Slides: {len(prs.slides)} | Done: {stats['done']}/{stats['total']} ({stats['done_pct']}%)")

# ============================================================
# BƯỚC 6: QA (chạy riêng sau khi save)
# ============================================================
# 1. Text QA:
#    extract-text /mnt/user-data/outputs/<file>.pptx | grep -iE "xx|lorem|TODO"
# 2. Visual QA:
#    python /mnt/skills/public/pptx/scripts/office/soffice.py \
#        --headless --convert-to pdf <file>.pptx
#    pdftoppm -jpeg -r 100 <file>.pdf slide
#    # Review từng slide-*.jpg
# 3. Fix lỗi overflow bằng cách rút gọn header (≤20 ký tự)
